from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BASE = "D:/Git/PianoHandSimulator/docs/중간보고서"
IMG  = lambda f: os.path.join(BASE, f)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 (화이트 테마) ─────────────────────
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
C_DARK   = RGBColor(0x2B, 0x2B, 0x2B)
C_MID    = RGBColor(0x55, 0x55, 0x55)
C_LIGHT  = RGBColor(0x99, 0x99, 0x99)
C_LINE   = RGBColor(0xCC, 0xCC, 0xCC)
C_BG2    = RGBColor(0xF5, 0xF5, 0xF5)   # 행 교대 배경
C_FIN    = RGBColor(0xC0, 0x5A, 0x00)   # Fingering 주황 (짙게)
C_IK     = RGBColor(0x00, 0x7A, 0x6E)   # IK 청록 (짙게)
C_SKIN   = RGBColor(0x6A, 0x35, 0x9C)   # Skinning 보라 (짙게)
C_BLUE   = RGBColor(0x1E, 0x5F, 0xA8)   # 중립 파랑

BLANK = prs.slide_layouts[6]

# ── 헬퍼 함수 ────────────────────────────────────
def bg(slide):
    r = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = C_WHITE
    r.line.fill.background()

def rect(slide, l, t, w, h, color, line_color=None):
    r = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color
    if line_color:
        r.line.color.rgb = line_color
    else:
        r.line.fill.background()
    return r

def txt(slide, text, l, t, w, h, size=14, bold=False, color=C_DARK,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.italic = italic
    return tb

def title_slide(slide, main_title, sub=None, accent=C_BLUE):
    """슬라이드 상단 타이틀 영역 (흰 배경, 하단 선)"""
    bg(slide)
    # 좌측 강조 바
    rect(slide, 0.5, 0.35, 0.06, 0.9, accent)
    txt(slide, main_title, 0.7, 0.3, 12.0, 0.7,
        size=28, bold=True, color=C_BLACK)
    if sub:
        txt(slide, sub, 0.7, 0.98, 12.0, 0.35,
            size=13, color=C_LIGHT, italic=True)
    # 구분선
    rect(slide, 0.5, 1.35, 12.33, 0.025, C_LINE)

def img(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        return
    if h:
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

def bullet_block(slide, lines, l, t, w, size=13, gap=0.4):
    """(indent, text, bold, color) 튜플 리스트를 세로로 출력"""
    for i, item in enumerate(lines):
        if isinstance(item, str):
            indent, text, bold, color = 0, item, False, C_DARK
        else:
            indent = item[0]; text = item[1]
            bold   = item[2] if len(item) > 2 else False
            color  = item[3] if len(item) > 3 else C_DARK
        txt(slide, text, l + indent * 0.35, t + i * gap, w - indent * 0.35,
            gap + 0.05, size=size, bold=bold, color=color)

def hline(slide, t, l=0.5, r=12.83):
    rect(slide, l, t, r - l, 0.02, C_LINE)

def section_label(slide, text, l, t, color=C_BLUE):
    txt(slide, text, l, t, 12.0, 0.45, size=13, bold=True, color=color)
    rect(slide, l, t + 0.42, 12.33 - l, 0.018, color)

# ═══════════════════════════════════════════════
# 슬라이드 1 — 표지
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)

# 상단 굵은 선
rect(sl, 0, 0, 13.33, 0.08, C_BLUE)

txt(sl, "피아노 연주 시뮬레이션 및 고정밀 렌더링 시스템",
    1.0, 2.0, 11.33, 1.1, size=34, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
txt(sl, "졸업 프로젝트 중간 보고서",
    1.0, 3.15, 11.33, 0.6, size=20, color=C_MID, align=PP_ALIGN.CENTER)

rect(sl, 4.0, 3.9, 5.33, 0.025, C_LINE)

info = [
    ("과목", "졸업프로젝트 3201 — 2팀"),
    ("팀원", "정근녕  ·  한승현  ·  이수민  ·  곽경민"),
    ("제출일", "2026년 5월 12일"),
]
for i, (k, v) in enumerate(info):
    txt(sl, k, 3.5, 4.15 + i * 0.55, 1.5, 0.5, size=13, color=C_LIGHT, align=PP_ALIGN.RIGHT)
    txt(sl, v, 5.2, 4.15 + i * 0.55, 5.0, 0.5, size=13, color=C_DARK)

# 하단 선
rect(sl, 0, 7.42, 13.33, 0.08, C_BLUE)

# ═══════════════════════════════════════════════
# 슬라이드 2 — 목차
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "목차")

items = [
    ("01", "프로젝트 개요"),
    ("02", "전체 시스템 아키텍처"),
    ("03", "진행 현황 요약"),
    ("04", "운지법 결정 엔진 (01_Fingering)"),
    ("05", "IK 기반 모션 생성 (02_IK)"),
    ("06", "커스텀 셰이더 파이프라인 (03_Skinning)"),
    ("07", "Chaos Flesh PBD Skinning (03_Skinning)"),
    ("08", "기술적 도전 및 해결"),
    ("09", "향후 계획 및 결론"),
]
for i, (num, title) in enumerate(items):
    row_y = 1.5 + i * 0.61
    if i % 2 == 1:
        rect(sl, 0.5, row_y, 12.33, 0.58, C_BG2)
    txt(sl, num, 0.6, row_y + 0.1, 0.7, 0.45, size=14, bold=True, color=C_BLUE)
    txt(sl, title, 1.4, row_y + 0.1, 11.0, 0.45, size=14, color=C_DARK)

# ═══════════════════════════════════════════════
# 슬라이드 3 — 프로젝트 개요
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "프로젝트 개요", "기획 배경 · 목표 · 파이프라인")

# 기획 배경
section_label(sl, "기획 배경", 0.5, 1.45)
probs = [
    (0, "기존 Linear Blend Skinning(LBS)의 한계"),
    (1, "관절 뭉개짐(Candy Wrapper Artifact), 주름·피하조직의 비사실적 표현"),
    (0, "MIDI → 실제 피아니스트 수준 운지법 자동 재현 도구의 부재"),
    (1, "기존 게임/시뮬레이션 엔진에서 해부학적 제약을 반영한 운지법 시스템 없음"),
]
bullet_block(sl, probs, 0.5, 1.95, 12.3, size=13, gap=0.42)

hline(sl, 4.1)

# 프로젝트 목표
section_label(sl, "프로젝트 목표", 0.5, 4.18)
goals = [
    (0, "MIDI-to-Motion:  MIDI 파싱 → 해부학적 운지법 자동 계산 → IK 애니메이션 데이터 생성", False, C_DARK),
    (0, "High-Fidelity Rendering:  PBD 조직 변형 + Tension Map 기반 동적 노멀로 주름·핏줄 표현", False, C_DARK),
]
bullet_block(sl, goals, 0.5, 4.65, 12.3, size=13, gap=0.5)

hline(sl, 5.75)

# 통합 파이프라인
section_label(sl, "통합 파이프라인", 0.5, 5.82)
stages = ["MIDI 입력", "운지법 계산", "IK 모션 생성", "피부 렌더링"]
colors = [C_MID, C_FIN, C_IK, C_SKIN]
for i, (label, color) in enumerate(zip(stages, colors)):
    x = 0.5 + i * 3.1
    rect(sl, x, 6.3, 2.8, 0.75, color)
    txt(sl, label, x, 6.35, 2.8, 0.65, size=14, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, "→", x + 2.8, 6.52, 0.3, 0.4, size=16, bold=True,
            color=C_MID, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# 슬라이드 4 — 팀 구성 & 기술 스택
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "팀 구성 및 기술 스택")

members = [
    ("정근녕", "01_Fingering", "MIDI 파서 및 DP 기반 운지법 결정 알고리즘", C_FIN),
    ("한승현", "02_IK",        "Jacobian 기반 역운동학(IK) 모션 생성",        C_IK),
    ("이수민", "03_Skinning",  "UE5 커스텀 셰이더, Tension Map, 카메라/UI",  C_SKIN),
    ("곽경민", "03_Skinning",  "Chaos Flesh PBD 피부 시뮬레이션, 핏줄·주름", C_SKIN),
]
for i, (name, mod, role, color) in enumerate(members):
    x = 0.5 + (i % 2) * 6.4
    y = 1.5 + (i // 2) * 1.8
    rect(sl, x, y, 5.9, 1.6, C_BG2, line_color=C_LINE)
    rect(sl, x, y, 0.06, 1.6, color)
    txt(sl, name, x + 0.2, y + 0.12, 2.5, 0.55, size=18, bold=True, color=C_BLACK)
    txt(sl, f"[{mod}]", x + 2.8, y + 0.15, 2.8, 0.45, size=12, color=color)
    txt(sl, role, x + 0.2, y + 0.75, 5.5, 0.75, size=12, color=C_MID)

hline(sl, 5.25)
section_label(sl, "기술 스택", 0.5, 5.32)
stacks = [
    ("Python / DP", C_FIN),
    ("C++ / Jacobian IK", C_IK),
    ("Unreal Engine 5.5", C_SKIN),
    ("Chaos Flesh / PBD", C_SKIN),
    ("HLSL / SVE / RDG", C_BLUE),
    ("MIDI / JSON", C_MID),
]
for i, (s, c) in enumerate(stacks):
    x = 0.5 + i * 2.05
    rect(sl, x, 5.82, 1.95, 0.52, C_BG2, line_color=c)
    txt(sl, s, x, 5.87, 1.95, 0.42, size=11, bold=True, color=c, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# 슬라이드 5 — 전체 시스템 아키텍처
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "전체 시스템 아키텍처")

# 입력 블록
rect(sl, 0.3, 2.8, 1.5, 0.8, C_BG2, line_color=C_MID)
txt(sl, "MIDI\n입력", 0.3, 2.75, 1.5, 0.9, size=13, bold=True,
    color=C_MID, align=PP_ALIGN.CENTER)

modules = [
    ("01_Fingering\n정근녕", ["MIDI 파서", "양손 분리", "V5 DP Solver", "JSON 출력"], C_FIN,  2.1),
    ("02_IK\n한승현",        ["IK 목표 설정", "Jacobian IK", "ROM 클램핑", "Bone Transform"], C_IK,  5.4),
    ("03_Skinning\n이수민·곽경민", ["Skeletal Mesh", "Chaos Flesh PBD", "커스텀 셰이더", "렌더링 출력"], C_SKIN, 8.7),
]
for (header, nodes, color, mx) in modules:
    rect(sl, mx, 1.5, 3.5, 4.5, C_BG2, line_color=color)
    rect(sl, mx, 1.5, 3.5, 0.5, color)
    txt(sl, header, mx, 1.52, 3.5, 0.46, size=11, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER)
    for j, n in enumerate(nodes):
        ny = 2.15 + j * 0.82
        rect(sl, mx + 0.2, ny, 3.1, 0.58, C_WHITE, line_color=color)
        txt(sl, n, mx + 0.2, ny + 0.05, 3.1, 0.48, size=12,
            color=color, align=PP_ALIGN.CENTER, bold=True)
        if j < 3:
            txt(sl, "↓", mx + 1.55, ny + 0.6, 0.4, 0.22,
                size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)
    txt(sl, "▶", mx - 0.32, 3.1, 0.32, 0.4, size=14, bold=True,
        color=C_MID, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# 슬라이드 6 — 진행 현황 요약
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "진행 현황 요약")

rows = [
    ("MIDI 파서",               "정근녕", "Format 0/1 파싱, 화음 그룹화, 양손 분리",                  "완료",           C_IK),
    ("운지법 엔진 V5",           "정근녕", "Polyphonic Voice Leading, 왼손 역전 수정, 시뮬레이터",     "완료",           C_IK),
    ("Jacobian IK",             "한승현", "DLS IK, ROM 제약, 베지어 손목 궤적, C/F 코드 검증",        "알고리즘 검증 완료", C_IK),
    ("UE5 IK 파이프라인",        "한승현", "UE5 Animation 시스템 데이터 전달",                        "진행 중",        C_BLUE),
    ("커스텀 셰이더 파이프라인",  "이수민", "SVE 삽입, RDG 자원 관리, HLSL 4-mode 분기, 스레드 동기화", "사전 검증 완료",  C_IK),
    ("물리 기반 Skinning 셰이더","이수민", "LBS→컴퓨트 셰이더 전환, 본 매트릭스 SBO 연동",            "진행 중",        C_BLUE),
    ("Chaos Flesh 에셋 구성",    "곽경민", "Dataflow 그래프, Capsule 콜리전 17개, 사면체 볼륨",        "에셋 구성 완료",  C_IK),
    ("Chaos Flesh BP 연동",      "곽경민", "캐릭터 블루프린트 Flesh 컴포넌트 연동",                   "버전 이슈로 중단", RGBColor(0xC0,0x20,0x20)),
]
hdrs  = ["모듈", "담당", "주요 완료 항목", "상태"]
col_w = [2.3,   1.0,    7.0,             2.0]
col_x = [0.5,   2.85,   3.9,             10.95]

for ci, (h, w, x) in enumerate(zip(hdrs, col_w, col_x)):
    rect(sl, x, 1.45, w, 0.42, C_BLACK)
    txt(sl, h, x + 0.05, 1.47, w - 0.1, 0.38,
        size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    y = 1.9 + ri * 0.57
    if ri % 2 == 1:
        rect(sl, 0.5, y, 12.45, 0.54, C_BG2)
    hline(sl, y + 0.54, l=0.5, r=12.95)
    for ci, (val, w, x) in enumerate(zip(row[:4], col_w, col_x)):
        clr = row[4] if ci == 3 else C_DARK
        bd  = ci == 3
        txt(sl, val, x + 0.05, y + 0.07, w - 0.1, 0.42,
            size=11, bold=bd, color=clr, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# 슬라이드 7 — 운지법 엔진 개요
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "운지법 결정 엔진", "01_Fingering — 정근녕", accent=C_FIN)

# 알고리즘 진화
section_label(sl, "알고리즘 진화 과정", 0.5, 1.45, color=C_FIN)
versions = [
    ("V1/V2", "개별 음표 단위 DP, 기초 비용 함수", "화음 처리 미흡, 손가락 꼬임 빈발"),
    ("V4",    "화음 그룹 DP, ROM 적용, MAX_SPAN 상수화", "Monophonic 가정 유지, 음악적 맥락 미반영"),
    ("V5 ★",  "Polyphonic Voice Leading, 3성부(MELODY/BASS/INNER) 분리", "현재 사용 버전"),
]
for i, (v, feat, note) in enumerate(versions):
    y = 1.95 + i * 1.42
    is_cur = v.startswith("V5")
    lc = C_FIN if is_cur else C_LINE
    rect(sl, 0.5, y, 5.9, 1.28, C_BG2 if not is_cur else RGBColor(0xFF,0xF3,0xE8), line_color=lc)
    txt(sl, v,    0.65, y + 0.08, 1.0, 0.45, size=15, bold=True, color=C_FIN)
    txt(sl, feat, 0.65, y + 0.48, 5.6, 0.4,  size=12, color=C_DARK)
    txt(sl, f"→ {note}", 0.65, y + 0.85, 5.6, 0.35, size=11, color=C_LIGHT, italic=True)

# Chord-based DP 구조
section_label(sl, "Chord-based DP 구조", 6.6, 1.45, color=C_FIN)
dp_lines = [
    (0, "30ms 이내 음표를 화음 그룹으로 묶어 처리"),
    (0, "State: 화음에 배정된 손가락 조합 (예: 1,3,5)"),
    (0, "Monotonicity 제약: 피치↑ = 손가락번호↑  → 꼬임 원천 차단"),
    (0, "MAX_SPAN: 손가락 쌍별 최대 건반 간격 상수화  → 가동 범위 초과 배정 차단"),
    (0, ""),
    (0, "비용 함수:  W_move + W_difficulty + W_anatomical + W_role", True, C_FIN),
]
bullet_block(sl, dp_lines, 6.6, 1.95, 6.2, size=12, gap=0.48)

# ═══════════════════════════════════════════════
# 슬라이드 8 — V5 Polyphonic Voice Leading
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "V5 Polyphonic Voice Leading", "01_Fingering — 성부 분리 · 손목 모델링 · 실시간 시뮬레이터", accent=C_FIN)

section_label(sl, "3성부 분리 구조", 0.5, 1.45, color=C_FIN)
voices = [
    ("MELODY", "4번, 5번 우선", "선율 Legato 확보, 엄지 사용 기피",     C_FIN),
    ("BASS",   "5번, 1번 우선", "저음역 안정 지지",                      RGBColor(0x1E,0x5F,0xA8)),
    ("INNER",  "2·3·4번",      "화음 충전, 나머지 음 효율 배치",         C_MID),
]
for i, (name, fingers, desc, color) in enumerate(voices):
    x = 0.5 + i * 4.25
    rect(sl, x, 1.95, 4.0, 1.5, C_BG2, line_color=color)
    rect(sl, x, 1.95, 4.0, 0.1, color)
    txt(sl, name,    x + 0.15, 2.12, 3.7, 0.45, size=16, bold=True, color=color)
    txt(sl, f"손가락: {fingers}", x + 0.15, 2.6,  3.7, 0.38, size=12, color=C_DARK)
    txt(sl, desc,    x + 0.15, 2.98, 3.7, 0.38, size=12, color=C_MID)

hline(sl, 3.6)
section_label(sl, "손목 물리 모델링", 0.5, 3.68, color=C_FIN)
wrist = [
    ("Wrist Yaw",  "손가락 건반 평균 위치 ↔ 손목 중심축 오프셋 계산",     "범위 ±35.0°"),
    ("Wrist Roll", "엄지·새끼가 흑건 누를 때 손목 기울어짐 모사",          "범위 ±20.0°"),
]
for i, (name, desc, rng) in enumerate(wrist):
    y = 4.18 + i * 0.62
    if i % 2 == 1: rect(sl, 0.5, y, 12.33, 0.58, C_BG2)
    txt(sl, name, 0.65, y + 0.1, 2.0, 0.4, size=13, bold=True, color=C_FIN)
    txt(sl, desc, 2.7,  y + 0.1, 7.5, 0.4, size=12, color=C_DARK)
    txt(sl, rng,  10.3, y + 0.1, 2.5, 0.4, size=12, color=C_MID, align=PP_ALIGN.RIGHT)

hline(sl, 5.5)
section_label(sl, "실시간 시뮬레이터 (Python / Pygame)", 0.5, 5.58, color=C_FIN)
sim = [
    (0, "이벤트 기반 MIDI 재생 — FPS 30 단위 동기화"),
    (0, "슬라이더 / 키보드 타임라인 탐색(Seek) 지원"),
    (0, "이중 마디 시각화: 안쪽 성부 역할 색상 / 끝단 손가락 번호 색상"),
]
bullet_block(sl, sim, 0.5, 6.05, 12.3, size=12, gap=0.42)

# ═══════════════════════════════════════════════
# 슬라이드 9 — 시뮬레이터 이미지
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "실시간 시뮬레이터 결과", "01_Fingering — Python/Pygame 시각화", accent=C_FIN)

img(sl, IMG("Fingering_ (1).png"), 0.5,  1.5, 5.9, 5.2)
img(sl, IMG("Fingering_ (2).png"), 6.9, 1.5, 5.9, 5.2)

hline(sl, 6.8)
txt(sl, "왼손(파랑)·오른손(빨강) 골격이 건반에 매핑",
    0.5, 6.88, 5.9, 0.45, size=11, color=C_MID, align=PP_ALIGN.CENTER)
txt(sl, "다중 손가락 동시 타건 + 하단 타임라인 슬라이더",
    6.9, 6.88, 5.9, 0.45, size=11, color=C_MID, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# 슬라이드 10 — IK 시스템
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "IK 기반 모션 생성", "02_IK — 한승현", accent=C_IK)

# 설계 방향
section_label(sl, "설계 방향", 0.5, 1.45, color=C_IK)
txt(sl, "오프라인 사전 계산 방식 채택 — 수렴 속도보다 정확도·제약 처리 유연성 우선\n"
        "DLS(Damped Least Squares) Jacobian IK 적용으로 특이점 발산 방지",
    0.5, 1.95, 6.0, 0.8, size=12, color=C_DARK)

section_label(sl, "DLS IK 연산 과정", 0.5, 2.85, color=C_IK)
steps = [
    "① 건반 좌표 ↔ End-effector 오차 벡터 산출",
    "② computeFK — 관절 글로벌 변환 갱신",
    "③ 회전축 × 거리벡터 외적 → 3×9 Jacobian 행렬 구성",
    "④ JJᵀ + λ²I 역행렬로 특이점 발산 방지  (λ = 0.05)",
    "⑤ tolerance 미만 수렴 시 반복 종료",
]
for i, s in enumerate(steps):
    y = 3.35 + i * 0.65
    if i % 2 == 0: rect(sl, 0.5, y, 5.9, 0.62, C_BG2)
    hline(sl, y + 0.62, l=0.5, r=6.4)
    rect(sl, 0.5, y, 0.05, 0.62, C_IK)
    txt(sl, s, 0.65, y + 0.1, 5.6, 0.42, size=12, color=C_DARK)

# ROM 표
section_label(sl, "관절 ROM 제약", 6.6, 1.45, color=C_IK)
rom_rows = [
    ("MCP",      "0°–90°",  "0°–20°", "±20°"),
    ("PIP",      "0°–100°", "0°–10°", "—"),
    ("DIP",      "0°–80°",  "0°–5°",  "—"),
    ("엄지 CMC", "0°–50°",  "0°–50°", "±40°"),
    ("엄지 MCP", "0°–60°",  "0°",     "—"),
    ("엄지 IP",  "0°–80°",  "0°–5°",  "—"),
]
hdrs2 = ["관절", "굴곡", "신전", "내외전"]
cw    = [1.5, 1.3, 1.3, 1.3]
cx    = [6.6, 8.15, 9.5, 10.85]
for ci, (h, w, x) in enumerate(zip(hdrs2, cw, cx)):
    rect(sl, x, 1.9, w, 0.4, C_BLACK)
    txt(sl, h, x, 1.92, w, 0.36, size=12, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER)
for ri, row in enumerate(rom_rows):
    y = 2.33 + ri * 0.55
    if ri % 2 == 1: rect(sl, 6.6, y, 5.85, 0.52, C_BG2)
    hline(sl, y + 0.52, l=6.6, r=12.45)
    for ci, (val, w, x) in enumerate(zip(row, cw, cx)):
        txt(sl, val, x + 0.05, y + 0.08, w - 0.1, 0.36,
            size=12, color=C_DARK, align=PP_ALIGN.CENTER)

# 베지어
hline(sl, 6.73)
txt(sl, "베지어 손목 궤적:  코드 전환 시 손목이 자연스럽게 들렸다 내려오는 2차 베지어 곡선 적용",
    6.6, 6.8, 6.3, 0.55, size=12, color=C_MID, italic=True)

# ═══════════════════════════════════════════════
# 슬라이드 11 — IK 시각화
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "IK 시각화 결과", "02_IK — Jacobian IK 손가락 골격 렌더링", accent=C_IK)

img(sl, IMG("IK_Image.png"), 1.5, 1.5, 10.3, 5.5)
hline(sl, 7.1)
txt(sl, "C++ 레벨에서 구현된 DLS Jacobian IK 결과 — 각 손가락이 목표 건반 좌표로 수렴",
    1.5, 7.15, 10.3, 0.35, size=11, color=C_MID, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# 슬라이드 12 — 커스텀 셰이더 파이프라인
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "커스텀 셰이더 파이프라인", "03_Skinning — 이수민", accent=C_SKIN)

txt(sl, "최종 목표: UE5 렌더링 파이프라인 내 물리 기반 GPU Skinning 셰이더 삽입",
    0.5, 1.45, 12.3, 0.42, size=13, bold=True, color=C_SKIN)
txt(sl, "→ 사전 검증 목적의 단계형 시각 처리(Staged Rendering) 테스트베드 선구축",
    0.5, 1.85, 12.3, 0.38, size=12, color=C_MID, italic=True)
hline(sl, 2.3)

items4 = [
    ("① SVE 커스텀 패스 삽입",
     "OnPostEngineInit 콜백으로 SVE 지연 생성\n"
     "SubscribeToPostProcessingPass → Tonemap 단계에 후크"),
    ("② 글로벌 셰이더 & RDG 디스패치",
     "FStagedRenderPS : FGlobalShader 상속\n"
     "AddFullscreenPass — PSO·뷰포트 자동화, RDG 자원 자동 회수"),
    ("③ HLSL 다중 모드 분기",
     "CVar(stage.Mode) 기반 4가지 렌더 경로\n"
     "0: 검정 / 1: Sobel 에지 / 2: 포스터라이즈 / 3: 패스스루"),
    ("④ 게임-렌더 스레드 동기화",
     "FindConsoleVariable→Set() 한 줄로 CVar 동기화\n"
     "1키 입력 → 4단계 즉각 순환 전환 확인 완료"),
]
for i, (title, desc) in enumerate(items4):
    x = 0.5  + (i % 2) * 6.45
    y = 2.42 + (i // 2) * 2.38
    rect(sl, x, y, 6.1, 2.22, C_BG2, line_color=C_LINE)
    rect(sl, x, y, 0.05, 2.22, C_SKIN)
    txt(sl, title, x + 0.18, y + 0.1,  5.75, 0.48, size=14, bold=True, color=C_SKIN)
    txt(sl, desc,  x + 0.18, y + 0.62, 5.75, 1.45, size=12, color=C_DARK)

# ═══════════════════════════════════════════════
# 슬라이드 13 — 셰이더 검증 결과 이미지
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "커스텀 셰이더 검증 결과", "03_Skinning — stage.Mode 순환 전환", accent=C_SKIN)

imgs3 = [
    ("UnrealPipelineTest.png",   "Stage 1 — Sobel 에지 검출"),
    ("UnrealPipelineTest_2.png", "Stage 2 — 탈색 + 포스터라이즈"),
    ("UnrealPipelineTest_3.png", "Stage 3 — 패스스루 (원본 복원)"),
]
for i, (fname, caption) in enumerate(imgs3):
    x = 0.35 + i * 4.32
    img(sl, IMG(fname), x, 1.5, 4.1, 5.3)
    hline(sl, 6.88, l=x, r=x + 4.1)
    txt(sl, caption, x, 6.95, 4.1, 0.45,
        size=11, color=C_MID, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# 슬라이드 14 — Chaos Flesh 개요 & 콜리전
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "Chaos Flesh PBD Skinning", "03_Skinning — 곽경민", accent=C_SKIN)

section_label(sl, "개요", 0.5, 1.45, color=C_SKIN)
txt(sl, "LBS의 Candy Wrapper Artifact 해소를 위해 UE5 FEM 기반 소프트바디 시뮬레이션 적용\n"
        "사면체(Tetrahedral) 내부 볼륨 생성 → 뼈 움직임에 따른 피부 자연 변형 구현",
    0.5, 1.95, 5.9, 0.85, size=12, color=C_DARK)

hline(sl, 2.9)
section_label(sl, "손가락 콜리전 설계", 0.5, 2.98, color=C_SKIN)
coll_lines = [
    (0, "손가락 마디별 콜리전 없으면 피부 관통 아티팩트 발생"),
    (0, "기존 단일 손바닥 바디 → 각 마디에 Capsule 콜리전 추가"),
    (0, "총 17개 관절 본에 적용"),
    (0, "Capsule 선택 이유: Sphere 대비 관절 방향 충돌 정확도 우수"),
]
bullet_block(sl, coll_lines, 0.5, 3.5, 5.9, size=12, gap=0.48)

img(sl, IMG("ChaosFlesh_ (4).png"), 6.5, 1.5, 6.4, 4.5)
txt(sl, "Physics Asset — 손가락 마디별 Capsule 콜리전 배치",
    6.5, 6.1, 6.4, 0.38, size=11, color=C_MID, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# 슬라이드 15 — Dataflow 그래프
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "Flesh Asset Dataflow 그래프", "03_Skinning — 두 그래프로 분리 구성", accent=C_SKIN)

txt(sl, "Dataflow ①  Geometry 생성 및 전처리",
    0.5, 1.45, 12.3, 0.42, size=13, bold=True, color=C_SKIN)
img(sl, IMG("ChaosFlesh_ (6).png"), 0.5, 1.9, 12.3, 2.1)

hline(sl, 4.1)
txt(sl, "Dataflow ②  시뮬레이션 제어 및 바인딩",
    0.5, 4.18, 12.3, 0.42, size=13, bold=True, color=C_SKIN)
img(sl, IMG("ChaosFlesh_ (7).png"), 0.5, 4.63, 12.3, 2.1)

hline(sl, 6.85)
txt(sl, "분리 이유:  Geometry 생성(고비용)은 사전 1회 수행 후 재사용 / 바인딩은 캐릭터별 독립 적용",
    0.5, 6.92, 12.3, 0.38, size=11, color=C_MID, italic=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# 슬라이드 16 — Flesh 시뮬레이션 결과
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "Flesh 시뮬레이션 결과", "03_Skinning — 에디터 프리뷰 비교", accent=C_SKIN)

cf_imgs = [
    ("ChaosFlesh_ (1).png", "T-포즈 기본 상태"),
    ("ChaosFlesh_ (5).png", "굴곡 — LBS 단독  (아티팩트 확인)"),
    ("ChaosFlesh_ (3).png", "굴곡 — Flesh 적용 후  (볼륨 보존)"),
]
for i, (fname, cap) in enumerate(cf_imgs):
    x = 0.35 + i * 4.32
    img(sl, IMG(fname), x, 1.5, 4.1, 5.1)
    hline(sl, 6.7, l=x, r=x + 4.1)
    txt(sl, cap, x, 6.77, 4.1, 0.45,
        size=11, color=C_MID, align=PP_ALIGN.CENTER)

# 이슈
rect(sl, 0.5, 7.22, 12.33, 0.22, RGBColor(0xFF,0xEE,0xEE), line_color=RGBColor(0xC0,0x20,0x20))
txt(sl, "현재 이슈:  UE 5.6 ↔ 5.5 Flesh 컴포넌트 API 변경으로 블루프린트 연동 중단 → UE 5.5 다운그레이드 검토 중",
    0.6, 7.25, 12.1, 0.3, size=11, color=RGBColor(0xC0,0x20,0x20))

# ═══════════════════════════════════════════════
# 슬라이드 17 — 기술적 도전 및 해결
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "주요 기술적 도전 및 해결")

challenges = [
    ("왼손 역전 문제",
     "오른손 중심 피치-손가락 매핑을 양손에 동일 적용 → 엄지(1번)·새끼(5번) X자 꼬임 발생",
     "왼손 손가락 조합 역순 배정 + 교차 판정 미러링으로 해결",
     C_FIN),
    ("Jacobian 특이점",
     "특정 자세에서 표준 역행렬이 수치 불안정하여 IK 계산값 발산",
     "DLS 역행렬  JJᵀ + λ²I  적용 (λ=0.05) + 스텝 크기 제한",
     C_IK),
    ("SVE 초기화 타이밍",
     "StartupModule에서 NewExtension 호출 시 GEngine == nullptr → 크래시",
     "OnPostEngineInit 콜백에 SVE 생성 이전·지연 처리로 해결",
     C_SKIN),
    ("Chaos Flesh 버전 이슈",
     "UE 5.5→5.6 업그레이드로 Flesh 전용 컴포넌트 API 변경 → BP 연동 중단",
     "UE 5.5 다운그레이드 검토 중, 현재 에디터 프리뷰 수준 동작 확인",
     RGBColor(0xC0,0x20,0x20)),
]
for i, (title, cause, fix, color) in enumerate(challenges):
    x = 0.5  + (i % 2) * 6.45
    y = 1.45 + (i // 2) * 2.95
    rect(sl, x, y, 6.1, 2.75, C_BG2, line_color=C_LINE)
    rect(sl, x, y, 0.05, 2.75, color)
    txt(sl, title, x + 0.18, y + 0.1,  5.75, 0.5, size=15, bold=True, color=color)
    txt(sl, "원인:",  x + 0.18, y + 0.7, 0.7, 0.38, size=11, bold=True, color=C_LIGHT)
    txt(sl, cause,   x + 0.95, y + 0.7, 4.98, 0.75, size=11, color=C_DARK)
    txt(sl, "해결:",  x + 0.18, y + 1.55, 0.7, 0.38, size=11, bold=True, color=C_LIGHT)
    txt(sl, fix,     x + 0.95, y + 1.55, 4.98, 0.75, size=11, color=C_DARK)

# ═══════════════════════════════════════════════
# 슬라이드 18 — 향후 계획 & 결론
# ═══════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
title_slide(sl, "향후 계획 및 결론")

plans = [
    ("01_Fingering", "복잡한 대위법 대응력 향상\nUE5 JSON 데이터 테이블 임포터 개발",            C_FIN),
    ("02_IK",        "UE5 MetaHuman 손 모델 IK 파이프라인 연동\n다중 손가락 충돌 방지 로직 통합", C_IK),
    ("03_Skinning (셰이더)", "PP 후크 → 정점 처리 단계 전환\nLBS/DQS 통합 컴퓨트 셰이더 구현",   C_SKIN),
    ("03_Skinning (Flesh)",  "UE 5.5 전환 후 BP 연동 완성\nTension Map 핏줄·주름 렌더링 구현",    C_SKIN),
]
for i, (mod, goal, color) in enumerate(plans):
    x = 0.5  + (i % 2) * 6.45
    y = 1.45 + (i // 2) * 2.0
    rect(sl, x, y, 6.1, 1.8, C_BG2, line_color=C_LINE)
    rect(sl, x, y, 0.05, 1.8, color)
    txt(sl, mod,  x + 0.2, y + 0.12, 5.7, 0.55, size=14, bold=True, color=color)
    txt(sl, goal, x + 0.2, y + 0.72, 5.7, 0.95, size=12, color=C_DARK)

hline(sl, 5.55)
section_label(sl, "결론", 0.5, 5.62, color=C_BLUE)
conclusions = [
    "01_Fingering:  V5 Polyphonic 엔진 완성 — 3성부 분리·ROM 기반 운지법 JSON 출력 안정화",
    "02_IK:  DLS Jacobian IK 알고리즘 검증 완료 — UE5 파이프라인 연동 진행 중",
    "03_Skinning (셰이더):  SVE·RDG·HLSL 4개 핵심 기술 사전 검증 완료",
    "03_Skinning (Flesh):  Dataflow 구성 완료, 버전 이슈 해소 후 블루프린트 연동 예정",
]
for i, c in enumerate(conclusions):
    txt(sl, c, 0.5, 6.12 + i * 0.36, 12.3, 0.35, size=12, color=C_DARK)

out = os.path.join(BASE, "중간보고서.pptx")
prs.save(out)
print(f"저장 완료: {out}")
